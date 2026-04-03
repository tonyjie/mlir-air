#!/usr/bin/env python3
"""Build LLAMA progress report PPTX from AMD template."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

TEMPLATE = "AMD_slide_template.pptx"
OUTPUT = "llama_report.pptx"

# AMD brand colors
AMD_RED = RGBColor(0xED, 0x1C, 0x24)
AMD_DARK = RGBColor(0x1A, 0x1A, 0x2E)
AMD_GRAY = RGBColor(0x58, 0x59, 0x5B)
AMD_GREEN = RGBColor(0x00, 0x96, 0x32)
AMD_LIGHT_GRAY = RGBColor(0xD9, 0xD9, 0xD9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

# Layout indices from template inspection
LY_TITLE = 0  # Title Slide - No Image
LY_CONTENT = 3  # Title and Content (bulleted)
LY_TEXT = 4  # Title & Non-bulleted text
LY_TWO_COL = 5  # Two Content
LY_DIVIDER = 26  # Divider slide
LY_CLOSING = 30  # Closing logo slide
LY_TITLE_ONLY = 7  # Title Only


def set_text(shape, text, font_size=None, bold=None, color=None, alignment=None):
    """Set simple text on a shape, clearing existing content."""
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    if alignment:
        p.alignment = alignment
    if font_size or bold is not None or color:
        run = p.runs[0]
        if font_size:
            run.font.size = Pt(font_size)
        if bold is not None:
            run.font.bold = bold
        if color:
            run.font.color.rgb = color


def add_bullet(tf, text, level=0, font_size=14, bold=False, color=None):
    """Add a bullet point to a text frame."""
    p = tf.add_paragraph()
    p.text = text
    p.level = level
    p.space_after = Pt(4)
    p.space_before = Pt(1)
    p.font.size = Pt(font_size)
    if text and p.runs:
        run = p.runs[0]
        run.font.size = Pt(font_size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color


def add_table(
    slide,
    rows_data,
    left,
    top,
    width,
    height,
    col_widths=None,
    header_color=AMD_DARK,
    font_size=10,
):
    """Add a formatted table to a slide."""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for row_idx, row_data in enumerate(rows_data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_text)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                if row_idx == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                else:
                    paragraph.font.color.rgb = BLACK

            # Header row styling
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF2, 0xF2, 0xF2)

    return table_shape


def add_callout_box(
    slide, text, left, top, width, height, fill_color=AMD_RED, font_size=18, bold=True
):
    """Add a colored callout/highlight box."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = text
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = WHITE
    return shape


# ============================================================
# Build presentation
# ============================================================

prs = Presentation(TEMPLATE)

# Remove existing template slides (keep just the master/layouts)
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].rId
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE])
title = slide.placeholders[0]
subtitle = slide.placeholders[12]
set_text(
    title, "LLAMA-3.2-1B BF16 Inference\non MLIR-AIR (NPU2)", font_size=32, bold=True
)
set_text(subtitle, "Progress Report  |  April 2026", font_size=14)


# ============================================================
# SLIDE 2: Results at a Glance
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_ONLY])
set_text(slide.placeholders[0], "Results at a Glance")

# Two big callout boxes
add_callout_box(
    slide,
    "Prefill: 1.92s\n30% faster than IRON",
    Inches(0.8),
    Inches(1.8),
    Inches(5.5),
    Inches(1.5),
    fill_color=AMD_RED,
    font_size=24,
)
add_callout_box(
    slide,
    "Decode: 351 ms/tok\n5% faster than IRON",
    Inches(7.0),
    Inches(1.8),
    Inches(5.5),
    Inches(1.5),
    fill_color=AMD_RED,
    font_size=24,
)

# Summary table below
add_table(
    slide,
    [
        ["Metric", "AIR (Ours)", "IRON (Ref)", "Improvement"],
        ["Prefill (seq=2048)", "1.92s", "2.744s", "30% faster"],
        ["Decode (steady-state)", "351 ms/tok", "370 ms/tok", "5% faster"],
        ["Correctness", 'Top-1 = "Paris"', "-", "corr 0.993 vs CPU"],
        ["Prefill inv/layer", "5", "~12", "58% fewer"],
        ["Decode inv/block", "10", "~12", "17% fewer"],
    ],
    Inches(0.8),
    Inches(3.7),
    Inches(11.7),
    Inches(2.8),
    col_widths=[Inches(3.0), Inches(2.8), Inches(2.8), Inches(3.1)],
    font_size=12,
)


# ============================================================
# SLIDE 3: LLAMA-3.2-1B Architecture
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[LY_CONTENT])
set_text(slide.placeholders[0], "Background: LLAMA-3.2-1B Architecture")

tf = slide.placeholders[10].text_frame
tf.clear()
tf.paragraphs[0].text = "1B parameter decoder-only transformer, 16 layers, BF16"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.bold = True

for item in [
    "Embedding dim = 2048, 32 Q heads / 8 KV heads (GQA), head dim = 64",
    "FFN hidden dim = 8192, vocab = 128,256",
    "Per-layer: RMSNorm -> Q/K/V -> RoPE -> Attention -> O + Add -> RMSNorm -> FFN -> Add",
]:
    add_bullet(tf, item, level=0, font_size=13)

add_bullet(tf, "", level=0, font_size=6)
add_bullet(tf, "Two inference phases:", level=0, font_size=14, bold=True)
add_bullet(
    tf,
    "Prefill: process full prompt (seq_len=2048) with GEMM (matrix-matrix)",
    level=1,
    font_size=13,
)
add_bullet(
    tf,
    "Decode: generate tokens one at a time with GEMV (matrix-vector)",
    level=1,
    font_size=13,
)

add_bullet(tf, "", level=0, font_size=6)
add_bullet(
    tf,
    "Reference: IRON (AMD's NPU framework) has working LLAMA at 2.744s prefill, 370ms/tok decode",
    level=0,
    font_size=13,
    bold=True,
    color=AMD_RED,
)


# ============================================================
# SLIDE 4: MLIR-AIR Starting Point
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[LY_CONTENT])
set_text(slide.placeholders[0], "MLIR-AIR Starting Point")

tf = slide.placeholders[10].text_frame
tf.clear()
tf.paragraphs[0].text = (
    "MLIR-AIR: compiler framework mapping AI workloads to AMD NPU via MLIR"
)
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.bold = True

add_bullet(
    tf,
    "Pipeline: Python IR gen -> AIR MLIR -> air-opt -> aircc -> aiecc -> xclbin/ELF",
    level=0,
    font_size=13,
)
add_bullet(
    tf,
    "Target: NPU2 (AIE2P, Strix) — 8 columns x 5 rows compute array, shared DDR",
    level=0,
    font_size=13,
)

add_bullet(tf, "", level=0, font_size=6)
add_bullet(tf, "What existed at project start:", level=0, font_size=14, bold=True)
add_bullet(
    tf,
    "Individual operator examples (GEMM, softmax, RMSNorm, etc.) at small sizes",
    level=1,
    font_size=13,
)
add_bullet(tf, "No end-to-end model inference", level=1, font_size=13)
add_bullet(tf, "No multi-launch ELF (operator fusion)", level=1, font_size=13)
add_bullet(tf, "No GEMV (decode) kernel", level=1, font_size=13)

add_bullet(tf, "", level=0, font_size=6)
add_bullet(tf, "What we built:", level=0, font_size=14, bold=True)
add_bullet(
    tf,
    "Full prefill pipeline: 7 fused kernel ELFs, 5 invocations/layer",
    level=1,
    font_size=13,
)
add_bullet(
    tf,
    "Full decode pipeline: 9 kernels, 10 invocations/block + CPU attention",
    level=1,
    font_size=13,
)
add_bullet(
    tf, "Multi-launch ELF infrastructure for operator fusion", level=1, font_size=13
)


# ============================================================
# SLIDE 5: Prefill Pipeline Architecture
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_ONLY])
set_text(slide.placeholders[0], "Prefill Pipeline: 5 Invocations per Layer")

add_table(
    slide,
    [
        ["#", "Operation", "Kernel", "Launches", "Time", "Key Detail"],
        ["1", "RMSNorm + Q/K/V GEMMs", "rms_attn_gemms", "4", "9 ms", "8-tile RMSNorm"],
        ["2", "RoPE Q+K", "rope_qk", "2 herds", "11 ms", "LUT-based rotation"],
        ["3", "Flash Attention GQA", "flash_attn", "1", "20 ms", "32Q/8KV heads"],
        ["4", "O GEMM + Residual Add", "o_proj_add", "2", "6 ms", "Fused residual"],
        ["5", "RMSNorm + FFN + Add", "ffn_full", "6", "52 ms", "8-tile RMSNorm"],
        ["", "Per-layer total", "", "5 calls", "~100 ms", ""],
        ["", "LM Head", "lm_head", "8", "171 ms", "8-partition ELF"],
        ["", "Total prefill", "", "", "1.92s", "IRON: 2.744s"],
    ],
    Inches(0.5),
    Inches(1.4),
    Inches(12.0),
    Inches(4.5),
    col_widths=[
        Inches(0.5),
        Inches(3.2),
        Inches(2.2),
        Inches(1.3),
        Inches(1.2),
        Inches(3.6),
    ],
    font_size=11,
)

# Annotation
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(6.1), Inches(12.0), Inches(0.6))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Each kernel = multi-launch ELF: multiple air.launch ops in one MLIR module, single xrt.run() per invocation"
p.font.size = Pt(11)
p.font.italic = True
p.font.color.rgb = AMD_GRAY


# ============================================================
# SLIDE 6: Prefill Optimization Journey
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_ONLY])
set_text(slide.placeholders[0], "Prefill Optimization: 18.67s -> 1.92s (10x)")

add_table(
    slide,
    [
        ["Step", "What Changed", "16-Layer Time", "vs IRON"],
        ["Baseline", "F32 scalar, no optimization", "18.67s", "7.8x slower"],
        ["+BF16 vectorize", "BF16 eltwise ops, vec=16", "13.40s", "5.6x"],
        ["+XRT reuse", "Cache XRT contexts", "8.77s", "3.7x"],
        ["+BO reuse", "Pre-allocate buffer objects", "6.49s", "2.7x"],
        ["+GEMM tuning", "Tile size optimization", "3.57s", "1.5x"],
        ["+FlashAttn NPU", "CPU -> NPU attention", "3.88s", "1.6x"],
        ["+Multi-launch", "FFN 4-launch, QKV 3-launch ELF", "2.45s", "0.92x"],
        ["+All merges", "5 inv/layer + bo.map() zero-copy", "1.81s", "0.74x"],
        ["+LM Head NPU", "8-partition ELF, static weight BOs", "2.05s*", "0.75x"],
        ["+8-tile RMSNorm", "Broadcast DMA fix, herd=[8,1]", "1.92s*", "0.70x"],
    ],
    Inches(0.5),
    Inches(1.3),
    Inches(12.0),
    Inches(5.2),
    col_widths=[Inches(2.2), Inches(5.2), Inches(2.0), Inches(2.6)],
    font_size=11,
)

txBox = slide.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(12.0), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "* Total prefill includes LM Head (171ms) + embedding overhead.  IRON total = 2.744s."
p.font.size = Pt(10)
p.font.italic = True
p.font.color.rgb = AMD_GRAY


# ============================================================
# SLIDE 7: Key Prefill Techniques
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[LY_CONTENT])
set_text(slide.placeholders[0], "Key Prefill Optimization Techniques")

tf = slide.placeholders[10].text_frame
tf.clear()

tf.paragraphs[0].text = "Multi-Launch ELF Fusion"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = AMD_RED
add_bullet(
    tf,
    "Stitch 2-6 air.launch ops into one MLIR module via text-based IR stitching",
    level=1,
    font_size=12,
)
add_bullet(
    tf,
    "Single xrt.run() per fused kernel — reduced 10 -> 5 invocations/layer",
    level=1,
    font_size=12,
)

add_bullet(
    tf, "bo.map() Zero-Copy Buffer I/O", level=0, font_size=14, bold=True, color=AMD_RED
)
add_bullet(
    tf,
    "Replace bo.read() (alloc+memcpy) with memory-mapped np.frombuffer(bo.map())",
    level=1,
    font_size=12,
)
add_bullet(
    tf,
    "Matching IRON's approach on Ryzen AI shared-memory architecture",
    level=1,
    font_size=12,
)

add_bullet(
    tf,
    "8-Tile RMSNorm with Broadcast Weight DMA",
    level=0,
    font_size=14,
    bold=True,
    color=AMD_RED,
)
add_bullet(
    tf,
    "Weight vector broadcast to all 8 tiles via single DMA (no tile-dependent offset)",
    level=1,
    font_size=12,
)
add_bullet(
    tf,
    "Required upstream compiler fix (stride=0 BD verifier). Standalone: 6ms -> 0.9ms (6.7x)",
    level=1,
    font_size=12,
)

add_bullet(
    tf,
    "NPU LM Head (8-Partition Multi-Launch)",
    level=0,
    font_size=14,
    bold=True,
    color=AMD_RED,
)
add_bullet(
    tf,
    "Vocab projection 2048x128K partitioned into 8 GEMMs, static weight BOs",
    level=1,
    font_size=12,
)
add_bullet(tf, "171ms vs IRON 217ms (21% faster) vs CPU 1,526ms", level=1, font_size=12)


# ============================================================
# SLIDE 8: Decode Pipeline
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_ONLY])
set_text(slide.placeholders[0], "Decode Pipeline: 10 Invocations per Block")

add_table(
    slide,
    [
        ["#", "Operation", "Kernel", "Herd", "Time"],
        ["1", "RMSNorm (pre-attn)", "rmsnorm", "[1,1]", "0.3 ms"],
        ["2", "Q+K+V GEMV (merged)", "qkv_gemv ELF", "[8,1]x3", "1.0 ms"],
        ["3-4", "RoPE Q, RoPE K", "rope_q/k", "[1,1]", "0.5 ms"],
        ["-", "CPU attention (GQA)", "numpy", "CPU", "~2 ms"],
        ["5", "O GEMV + Add (merged)", "o_gemv_add ELF", "[8,1]+[8,1]", "0.6 ms"],
        ["6", "RMSNorm (pre-FFN)", "rmsnorm", "[1,1]", "0.3 ms"],
        ["7", "Gate+Up GEMV (merged)", "gate_up_gemv ELF", "[8,1]x2", "2.5 ms"],
        ["8", "SiLU x mul", "silu_mul", "[8,1]", "0.3 ms"],
        ["9", "Down GEMV", "gemv_down", "[8,1]", "2.1 ms"],
        ["10", "Residual Add", "add", "[8,1]", "0.3 ms"],
        ["", "Per-block total", "", "10 calls", "~8 ms"],
    ],
    Inches(0.5),
    Inches(1.3),
    Inches(12.0),
    Inches(5.5),
    col_widths=[Inches(0.8), Inches(3.5), Inches(2.8), Inches(2.3), Inches(2.6)],
    font_size=11,
)


# ============================================================
# SLIDE 9: Decode Optimization Phases
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[LY_CONTENT])
set_text(slide.placeholders[0], "Decode: From Scratch to Faster than IRON")

tf = slide.placeholders[10].text_frame
tf.clear()
tf.paragraphs[0].text = (
    "Built GEMV kernel (matrix-vector) from scratch — didn't exist in MLIR-AIR"
)
tf.paragraphs[0].font.size = Pt(13)
tf.paragraphs[0].font.bold = True

add_bullet(tf, "", level=0, font_size=4)
add_bullet(
    tf,
    "Phase 1: First Working Pipeline  (~500 ms/tok)",
    level=0,
    font_size=14,
    bold=True,
    color=AMD_RED,
)
add_bullet(
    tf,
    "15 NPU calls/block, CPU SiLU, no BO caching, Python invoker overhead dominated",
    level=1,
    font_size=12,
)

add_bullet(
    tf,
    "Phase 2: Static Weight BOs + bo.map()  (~340 ms/tok)",
    level=0,
    font_size=14,
    bold=True,
    color=AMD_RED,
)
add_bullet(
    tf,
    "Per-layer BO isolation via bo_key (128 BO sets, 8 XRT contexts)",
    level=1,
    font_size=12,
)
add_bullet(
    tf,
    "Weights written once on first token, skipped on subsequent — saved 160 ms/tok",
    level=1,
    font_size=12,
)

add_bullet(
    tf,
    "Phase 3: Multi-Launch + NPU SiLU  (~351 ms/tok)",
    level=0,
    font_size=14,
    bold=True,
    color=AMD_RED,
)
add_bullet(
    tf,
    "Merged Q+K+V (3->1), O+Add (2->1), Gate+Up (2->1) into multi-launch ELFs",
    level=1,
    font_size=12,
)
add_bullet(
    tf,
    "Moved SiLU x mul to NPU, upgraded eltwise_add to [8,1] herd",
    level=1,
    font_size=12,
)

add_bullet(tf, "", level=0, font_size=4)
add_bullet(
    tf,
    "Steady-state: 351 ms/tok (IRON: 370 ms/tok) — AIR 5% faster",
    level=0,
    font_size=14,
    bold=True,
    color=AMD_GREEN,
)


# ============================================================
# SLIDE 10: AIR vs IRON Comparison
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_ONLY])
set_text(slide.placeholders[0], "AIR vs IRON: Architecture Comparison")

add_table(
    slide,
    [
        ["Aspect", "AIR (Ours)", "IRON (Reference)"],
        [
            "Kernel format",
            "Multi-launch ELF\n(2-8 air.launch per ELF)",
            "Runlist\n(multiple xclbin entries)",
        ],
        ["Operator fusion", "Text-based MLIR stitching", "Built-in runlist support"],
        [
            "Data path",
            "DDR -> L2 -> L1\n(MemTile staging)",
            "DDR -> L1 direct\n(ObjectFIFO)",
        ],
        ["Buffer I/O", "bo.map() zero-copy", "bo.map() zero-copy"],
        ["RMSNorm", "8-tile broadcast (0.9 ms)", "16-tile ObjectFIFO (4.3 ms)"],
        ["FFN (decode)", "4 separate NPU calls", "1 fused SwiGLU (5 entries)"],
        ["LM Head (prefill)", "8-launch ELF (171 ms)", "Single op (217 ms)"],
        ["Prefill dispatches/layer", "5", "~12"],
        ["Decode dispatches/block", "10", "~12"],
    ],
    Inches(0.5),
    Inches(1.3),
    Inches(12.0),
    Inches(5.2),
    col_widths=[Inches(3.2), Inches(4.4), Inches(4.4)],
    font_size=11,
)


# ============================================================
# SLIDE 11: Compiler Bugs
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[LY_CONTENT])
set_text(slide.placeholders[0], "Compiler Challenges Encountered")

tf = slide.placeholders[10].text_frame
tf.clear()
tf.paragraphs[0].text = "Bugs discovered and resolved during development:"
tf.paragraphs[0].font.size = Pt(13)

add_bullet(tf, "", level=0, font_size=4)
add_bullet(
    tf,
    "Broadcast DMA stride=0 rejection  [FIXED upstream]",
    level=0,
    font_size=13,
    bold=True,
    color=AMD_GREEN,
)
add_bullet(
    tf,
    "Multi-tile kernels with shared weights failed at aie.dma_bd verifier",
    level=1,
    font_size=12,
)
add_bullet(
    tf, "Fix enabled 8-tile RMSNorm: 6ms -> 0.9ms (6.7x faster)", level=1, font_size=12
)

add_bullet(
    tf,
    "Bare herd in multi-launch ELF  [WORKAROUND]",
    level=0,
    font_size=13,
    bold=True,
    color=RGBColor(0xFF, 0x99, 0x00),
)
add_bullet(
    tf,
    "Herd without air.segment silently dropped during legalization",
    level=1,
    font_size=12,
)
add_bullet(
    tf,
    "Solution: _wrap_ir_in_launch() adds segment wrapper automatically",
    level=1,
    font_size=12,
)

add_bullet(
    tf,
    "BF16 DMA stride limitation  [HARDWARE]",
    level=0,
    font_size=13,
    bold=True,
    color=AMD_RED,
)
add_bullet(
    tf, "AIE DMA requires innermost stride=1 for sub-32b types", level=1, font_size=12
)
add_bullet(
    tf,
    "Blocks DMA-only transpose for BF16 (would save 2 inv/layer in prefill)",
    level=1,
    font_size=12,
)

add_bullet(
    tf,
    "linalg_fill type mismatch  [OPEN]",
    level=0,
    font_size=13,
    bold=True,
    color=AMD_RED,
)
add_bullet(
    tf,
    "Different GEMV tile_m sizes produce incompatible func signatures in same module",
    level=1,
    font_size=12,
)
add_bullet(
    tf, "Blocks FFN full merge for decode (5 launches in 1 ELF)", level=1, font_size=12
)


# ============================================================
# SLIDE 12: Methodology
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[LY_CONTENT])
set_text(slide.placeholders[0], "Methodology: Documentation-Driven Development")

tf = slide.placeholders[10].text_frame
tf.clear()
tf.paragraphs[0].text = "Structured approach across multiple working sessions:"
tf.paragraphs[0].font.size = Pt(13)

add_bullet(tf, "", level=0, font_size=4)
add_bullet(
    tf,
    "Phased roadmap with clear entry/exit criteria",
    level=0,
    font_size=14,
    bold=True,
    color=AMD_RED,
)
add_bullet(
    tf,
    "Infrastructure -> Kernel validation -> Single layer -> Full model -> Performance",
    level=1,
    font_size=12,
)
add_bullet(
    tf,
    "Correctness verified before each performance optimization step",
    level=1,
    font_size=12,
)

add_bullet(
    tf,
    "Living documentation (8+ docs maintained throughout)",
    level=0,
    font_size=14,
    bold=True,
    color=AMD_RED,
)
add_bullet(
    tf,
    "LLAMA_PLAN.md: single-page status overview for session restoration",
    level=1,
    font_size=12,
)
add_bullet(
    tf,
    "perf_opt_prefill.md: optimization timeline with per-step measurements",
    level=1,
    font_size=12,
)
add_bullet(
    tf,
    "Per-kernel analysis docs and compiler bug reports with reproducers",
    level=1,
    font_size=12,
)

add_bullet(
    tf,
    "Claude Code as development partner",
    level=0,
    font_size=14,
    bold=True,
    color=AMD_RED,
)
add_bullet(
    tf,
    "Plan mode for design, subagents for parallel exploration, auto-memory for context",
    level=1,
    font_size=12,
)
add_bullet(
    tf,
    "Iterative loop: change -> compile -> run -> profile -> compare -> decide next",
    level=1,
    font_size=12,
)
add_bullet(
    tf,
    "Documentation-first: each session starts by reading LLAMA_PLAN.md",
    level=1,
    font_size=12,
)


# ============================================================
# SLIDE 13: Remaining Work
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_ONLY])
set_text(slide.placeholders[0], "Remaining Work & Future Directions")

add_table(
    slide,
    [
        ["Priority", "Task", "Expected Impact", "Status"],
        [
            "1",
            "Variable-length input sequences",
            "Usability (any prompt length)",
            "To investigate",
        ],
        ["2", "NPU LM Head for decode", "~40 ms/tok faster decode", "Unstarted"],
        [
            "3",
            "FFN full merge for decode",
            "~10 ms/tok (3 fewer dispatches)",
            "Blocked (type mismatch)",
        ],
        ["4", "NPU prefill for KV cache", "16s -> 2s init time", "Future"],
        ["5", "Unified prefill + decode script", "User-facing convenience", "Future"],
        ["6", "DMA transpose (prefill)", "5 -> 3 inv/layer", "Blocked (BF16 stride)"],
    ],
    Inches(0.5),
    Inches(1.5),
    Inches(12.0),
    Inches(3.5),
    col_widths=[Inches(1.0), Inches(4.0), Inches(4.0), Inches(3.0)],
    font_size=12,
)

txBox = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.0), Inches(1.2))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Variable-length input: Currently our pipeline requires fixed seq_len=2048. Short prompts are padded with EOS tokens, wasting compute. IRON handles arbitrary prompt lengths — investigating their approach is a key next step for practical usability."
p.font.size = Pt(12)
p.font.color.rgb = AMD_GRAY


# ============================================================
# SLIDE 14: Summary
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_ONLY])
set_text(slide.placeholders[0], "Summary")

add_callout_box(
    slide,
    "Prefill: 1.92s  (30% faster than IRON)",
    Inches(1.0),
    Inches(1.6),
    Inches(5.0),
    Inches(1.0),
    fill_color=AMD_RED,
    font_size=20,
)
add_callout_box(
    slide,
    "Decode: 351 ms/tok  (5% faster than IRON)",
    Inches(7.2),
    Inches(1.6),
    Inches(5.0),
    Inches(1.0),
    fill_color=AMD_RED,
    font_size=20,
)

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(3.0), Inches(11.5), Inches(3.5))
tf = txBox.text_frame
tf.word_wrap = True

items = [
    ("Built end-to-end LLAMA-3.2-1B inference on MLIR-AIR from scratch", True),
    ("Prefill: 7 fused kernel ELFs, 5 invocations/layer, NPU LM Head", False),
    ("Decode: 9 kernels, 10 invocations/block, static weight BO caching", False),
    ("", False),
    ("Key techniques that enabled this:", True),
    ("Multi-launch ELF fusion (text-based MLIR stitching)", False),
    ("bo.map() zero-copy buffer I/O", False),
    ("8-tile RMSNorm with broadcast weight DMA (compiler bug fixed)", False),
    ("Static weight BO caching with per-layer isolation", False),
    ("", False),
    ("Demonstrated MLIR-AIR can match and exceed IRON's performance", True),
]
for i, (text, is_bold) in enumerate(items):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(13)
    p.font.bold = is_bold
    p.space_after = Pt(3)
    if not text:
        p.space_after = Pt(1)
    if is_bold and text:
        p.font.color.rgb = AMD_DARK


# ============================================================
# SLIDE 15: Closing
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[LY_CLOSING])


# ============================================================
# Save
# ============================================================
output_path = (
    f"/home/jiajli/apps/mlir-air/programming_examples/llama3/docs/report/{OUTPUT}"
)
prs.save(output_path)
print(f"Saved to {output_path}")
print(f"Total slides: {len(prs.slides)}")
